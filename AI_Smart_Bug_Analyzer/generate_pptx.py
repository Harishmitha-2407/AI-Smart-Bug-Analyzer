"""
generate_pptx.py
Generates the AI Smart Bug Analyzer presentation (.pptx)
Run: python generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x1B, 0x3E)   # deep navy
BLUE      = RGBColor(0x15, 0x5F, 0xA4)   # primary blue
LIGHTBLUE = RGBColor(0x4F, 0xC3, 0xF7)   # accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xF0, 0xF5, 0xFF)
GRAY      = RGBColor(0x55, 0x66, 0x77)
GREEN     = RGBColor(0x1B, 0x8A, 0x4E)
ORANGE    = RGBColor(0xE6, 0x7E, 0x22)
DARKGRAY  = RGBColor(0x2C, 0x3E, 0x50)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # completely blank

# ══════════════════════════════════════════════════════════════════════════════
# Helper functions
# ══════════════════════════════════════════════════════════════════════════════

def add_slide():
    return prs.slides.add_slide(blank_layout)

def rect(slide, l, t, w, h, fill=None, line=None, radius=False):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, l, t, w, h,
          size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def add_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text

def header_bar(slide, title, subtitle=None):
    """Standard top navy bar with title text."""
    rect(slide, 0, 0, 13.33, 1.35, fill=NAVY)
    rect(slide, 0, 1.35, 13.33, 0.06, fill=LIGHTBLUE)
    txbox(slide, title, 0.4, 0.15, 11.5, 0.75,
          size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(slide, subtitle, 0.4, 0.82, 11, 0.44,
              size=14, color=LIGHTBLUE, align=PP_ALIGN.LEFT)

def footer_bar(slide, text="AI Smart Bug Analyzer & Fix Advisor  |  Phase 1"):
    rect(slide, 0, 7.15, 13.33, 0.35, fill=NAVY)
    txbox(slide, text, 0.3, 7.17, 12, 0.28,
          size=9, color=LIGHTBLUE, align=PP_ALIGN.LEFT)

def bullet_block(slide, items, l, t, w, h,
                 icon="▸", size=14, color=DARKGRAY, spacing=0.38):
    y = t
    for item in items:
        txbox(slide, f"{icon}  {item}", l, y, w, spacing + 0.04,
              size=size, color=color)
        y += spacing

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ══════════════════════════════════════════════════════════════════════════════
s1 = add_slide()

# Full navy background
rect(s1, 0, 0, 13.33, 7.5, fill=NAVY)
# Accent stripe
rect(s1, 0, 4.9, 13.33, 0.08, fill=LIGHTBLUE)
# Bottom accent
rect(s1, 0, 6.8, 13.33, 0.7, fill=RGBColor(0x0A, 0x12, 0x2A))

# Bug icon block
rect(s1, 0.5, 0.5, 0.9, 0.9, fill=LIGHTBLUE)
txbox(s1, "🐛", 0.52, 0.48, 0.86, 0.86, size=32, align=PP_ALIGN.CENTER)

# Title
txbox(s1, "AI Smart Bug Analyzer", 1.55, 0.45, 10.8, 0.7,
      size=38, bold=True, color=WHITE)
txbox(s1, "& Fix Advisor", 1.55, 1.1, 10.8, 0.6,
      size=32, bold=True, color=LIGHTBLUE)

# Subtitle tag
rect(s1, 1.55, 1.78, 4.2, 0.42, fill=BLUE)
txbox(s1, "Phase 1  ·  Tasks 1 – 5  ·  College Mini Project",
      1.65, 1.82, 4.0, 0.35, size=11, color=WHITE)

# Divider
rect(s1, 1.55, 2.4, 9.8, 0.04, fill=LIGHTBLUE)

# Team info grid
info = [
    ("👥  Team Members",   "Your Name 1  ·  Your Name 2  ·  Your Name 3"),
    ("🎓  Guide",          "Prof. Guide Name"),
    ("🏛️  Department",     "Department of Computer Science & Engineering"),
    ("🏫  College",        "Your College Name"),
    ("📅  Date",           "July 2026"),
]
y = 2.55
for label, value in info:
    txbox(s1, label, 1.55, y, 3.5, 0.34, size=11, color=LIGHTBLUE, bold=True)
    txbox(s1, value, 4.8,  y, 7.0, 0.34, size=11, color=WHITE)
    y += 0.38

add_notes(s1, "Welcome slide. Introduce the project title, team, and guide. "
              "Briefly mention this is a Phase 1 college mini project focused on "
              "building the foundation for an AI-powered bug analysis system.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
s2 = add_slide()
rect(s2, 0, 0, 13.33, 7.5, fill=OFFWHITE)
header_bar(s2, "Problem Statement", "Why do we need an intelligent bug management system?")
footer_bar(s2)

# Left pain-point cards
cards = [
    (BLUE,   "📋", "Manual Bug Tracking",    "Developers spend hours manually categorising and tracking bugs across spreadsheets and email threads."),
    (NAVY,   "🔍", "No Semantic Search",     "Keyword searches miss bugs described with different words, causing the same issue to be solved repeatedly."),
    (RGBColor(0x8B,0x00,0x00), "⚠️", "Inconsistent Reporting", "Bug reports lack structure — missing stack traces, OS info, or expected vs actual behaviour."),
    (RGBColor(0x4A,0x23,0x5A), "⏱️", "Wasted Developer Time",  "Without historical context, developers re-investigate known issues instead of applying proven fixes."),
]

x_positions = [0.35, 3.55, 6.75, 9.95]
for i, (col, icon, title, desc) in enumerate(cards):
    x = x_positions[i]
    rect(s2, x, 1.6, 3.0, 4.8, fill=col)
    txbox(s2, icon,  x+0.1, 1.7,  2.8, 0.7,  size=30, align=PP_ALIGN.CENTER)
    txbox(s2, title, x+0.1, 2.45, 2.8, 0.55, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s2, x+0.1, 3.08, 2.8, 0.04, fill=LIGHTBLUE)
    txbox(s2, desc,  x+0.15, 3.2, 2.7, 2.8,  size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_notes(s2, "Explain each pain point to the audience. Emphasise that manual bug "
              "tracking is error-prone and slow. The goal of this project is to build a "
              "structured, AI-ready system that addresses all four problems.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Project Objectives
# ══════════════════════════════════════════════════════════════════════════════
s3 = add_slide()
rect(s3, 0, 0, 13.33, 7.5, fill=OFFWHITE)
header_bar(s3, "Project Objectives", "What this project aims to achieve")
footer_bar(s3)

objectives = [
    ("01", BLUE,                        "Accept Bug Descriptions",       "Allow developers to paste any error message, exception, or bug narrative directly into the UI."),
    ("02", RGBColor(0x0A,0x6A,0x3F),    "Upload Log Files",              "Support .txt and .log file uploads so raw application logs are captured alongside descriptions."),
    ("03", RGBColor(0x6A,0x0A,0x6A),    "Modular Agent Architecture",    "Design five independent AI agents that each own a single responsibility in the analysis pipeline."),
    ("04", RGBColor(0x8B,0x45,0x00),    "Structured Knowledge Base",     "Define a rich CSV schema (14 fields) that stores bugs in a format ready for future vector indexing."),
    ("05", RGBColor(0x1A,0x3A,0x6A),    "AI-Ready Foundation",           "Build placeholder agents and clear interfaces so Phase 2 AI features can be dropped in without refactoring."),
]

y = 1.6
for num, col, title, desc in objectives:
    rect(s3, 0.4, y, 0.65, 0.65, fill=col)
    txbox(s3, num, 0.4, y+0.07, 0.65, 0.5, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s3, 1.15, y, 11.7, 0.65, fill=WHITE, line=RGBColor(0xCC,0xDD,0xEE))
    txbox(s3, title, 1.28, y+0.05, 3.5, 0.32, size=13, bold=True, color=NAVY)
    txbox(s3, desc,  4.9,  y+0.08, 7.8, 0.5,  size=11, color=GRAY)
    y += 0.84

add_notes(s3, "Walk through each objective. Highlight that objective 5 is about future-proofing — "
              "the agents are designed with clear interfaces so Phase 2 can plug in embeddings and "
              "LLM calls without touching the UI or submission logic.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Project Scope
# ══════════════════════════════════════════════════════════════════════════════
s4 = add_slide()
rect(s4, 0, 0, 13.33, 7.5, fill=OFFWHITE)
header_bar(s4, "Project Scope", "Phase 1 — Completed deliverables")
footer_bar(s4)

# Completed column
rect(s4, 0.4, 1.6, 5.9, 5.2, fill=WHITE, line=RGBColor(0xBB,0xCC,0xDD))
rect(s4, 0.4, 1.6, 5.9, 0.55, fill=GREEN)
txbox(s4, "✅  COMPLETED — Phase 1", 0.55, 1.68, 5.6, 0.4,
      size=13, bold=True, color=WHITE)

done = [
    "📖  Workflow & pipeline study",
    "🏗️  System architecture design",
    "🤖  Five AI agent class stubs",
    "🗄️  Knowledge base schema (14 fields)",
    "📋  Sample dataset — 20 realistic bugs",
    "🖥️  Streamlit submission UI",
    "✔️  Input validation & error handling",
    "💾  CSV persistence & file uploads",
]
y = 2.3
for item in done:
    txbox(s4, item, 0.65, y, 5.4, 0.38, size=12, color=DARKGRAY)
    y += 0.42

# Out of scope column
rect(s4, 7.0, 1.6, 5.9, 5.2, fill=WHITE, line=RGBColor(0xBB,0xCC,0xDD))
rect(s4, 7.0, 1.6, 5.9, 0.55, fill=RGBColor(0x4A,0x4A,0x6A))
txbox(s4, "🔒  PHASE 2 — Future Work", 7.15, 1.68, 5.6, 0.4,
      size=13, bold=True, color=WHITE)

future = [
    "🧠  Text embedding generation",
    "📦  ChromaDB vector store setup",
    "🔍  Semantic similarity search",
    "🤖  RAG pipeline implementation",
    "💡  Gemini AI fix recommendation",
    "📊  Similarity confidence scoring",
    "🔗  GitHub / Jira integration",
]
y = 2.3
for item in future:
    txbox(s4, item, 7.25, y, 5.4, 0.38, size=12, color=GRAY)
    y += 0.42

add_notes(s4, "Clearly separate what is done from what is planned. "
              "This scope boundary is important for the evaluators to understand "
              "that Phase 1 is a complete, working system — not a half-finished project.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Workflow Diagram
# ══════════════════════════════════════════════════════════════════════════════
s5 = add_slide()
rect(s5, 0, 0, 13.33, 7.5, fill=OFFWHITE)
header_bar(s5, "Workflow Diagram", "End-to-end bug submission pipeline")
footer_bar(s5)

# Flowchart nodes — centered at x=4.4
flow_nodes = [
    (BLUE,                        "👤  Developer",               "Bug description / log file"),
    (BLUE,                        "📥  Bug Submission Module",   "Accepts text & uploaded files"),
    (RGBColor(0x0A,0x6A,0x3F),    "✅  Validation",              "Checks type, size, empty inputs"),
    (RGBColor(0x6A,0x0A,0x6A),    "🔧  Preprocessing",          "Clean, lowercase, strip HTML"),
    (RGBColor(0x8B,0x45,0x00),    "💾  Temporary Storage",       "Saved to bug_dataset.csv"),
    (RGBColor(0x3A,0x3A,0x5A),    "🔒  Future AI Processing",    "Embedding → Retrieval → Gemini"),
]

box_w, box_h = 4.6, 0.7
x_center = 4.35
y_start  = 1.55
gap      = 0.3

for i, (col, label, sub) in enumerate(flow_nodes):
    y = y_start + i * (box_h + gap)
    rect(s5, x_center, y, box_w, box_h, fill=col)
    txbox(s5, label, x_center + 0.15, y + 0.05, box_w - 0.3, 0.38,
          size=13, bold=True, color=WHITE)
    txbox(s5, sub,   x_center + 0.15, y + 0.38, box_w - 0.3, 0.28,
          size=10, color=RGBColor(0xCC,0xDD,0xFF))
    # Arrow (except after last)
    if i < len(flow_nodes) - 1:
        arrow_y = y + box_h
        rect(s5, x_center + box_w/2 - 0.05, arrow_y, 0.1, gap, fill=LIGHTBLUE)
        # Arrowhead triangle approximation
        txbox(s5, "▼", x_center + box_w/2 - 0.18, arrow_y + gap - 0.2,
              0.36, 0.28, size=12, color=LIGHTBLUE, align=PP_ALIGN.CENTER)

# Right annotation panel
rect(s5, 9.6, 1.55, 3.3, 5.6, fill=WHITE, line=RGBColor(0xBB,0xCC,0xDD))
txbox(s5, "📌  Phase 1 Details", 9.75, 1.65, 2.9, 0.4, size=12, bold=True, color=NAVY)
notes_items = [
    "Text area input",
    "File uploader (.txt/.log)",
    "Character counter",
    "Size & type validation",
    "UUID bug ID generation",
    "UTC timestamp",
    "CSV row append",
    "Upload file persistence",
]
y2 = 2.15
for item in notes_items:
    txbox(s5, f"  •  {item}", 9.75, y2, 2.9, 0.35, size=10, color=GRAY)
    y2 += 0.36

add_notes(s5, "Walk through each step of the flowchart from top to bottom. "
              "Emphasise that the bottom node (Future AI Processing) is a placeholder "
              "only — it is not implemented in Phase 1. The dashed appearance signals this.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — System Architecture
# ══════════════════════════════════════════════════════════════════════════════
s6 = add_slide()
rect(s6, 0, 0, 13.33, 7.5, fill=OFFWHITE)
header_bar(s6, "System Architecture", "Component overview — Phase 1 active  ·  Phase 2 planned")
footer_bar(s6)

# Layer definitions: (y, height, bg_color, label, components)
layers = [
    (1.55, 0.9,  BLUE,                       "🖥️  User Interface  (Streamlit)",
     ["Submit Bug Tab", "Dataset Viewer Tab", "Documentation Tab"]),

    (2.65, 0.9,  RGBColor(0x0A,0x5A,0x3A),   "🤖  Agent Layer — Phase 1",
     ["SubmissionAgent", "PreprocessingAgent"]),

    (3.75, 0.9,  RGBColor(0x8B,0x45,0x00),   "🔧  Utilities",
     ["Validator", "Cleaner", "FileHandler", "Logger"]),

    (4.85, 0.9,  RGBColor(0x35,0x1C,0x75),   "💾  Storage — Phase 1",
     ["bug_dataset.csv", "uploaded_logs/"]),

    (5.95, 0.7,  RGBColor(0x3A,0x3A,0x5A),   "🔒  Phase 2 Planned",
     ["EmbeddingAgent", "ChromaDB", "RetrievalAgent", "FixAdvisorAgent (Gemini)"]),
]

for ly, lh, col, label, comps in layers:
    rect(s6, 0.35, ly, 12.6, lh, fill=col)
    txbox(s6, label, 0.5, ly + 0.08, 3.5, 0.4, size=12, bold=True, color=WHITE)
    cx = 4.2
    for comp in comps:
        rect(s6, cx, ly + 0.14, 1.9, lh - 0.28, fill=WHITE)
        txbox(s6, comp, cx + 0.05, ly + lh/2 - 0.15, 1.8, 0.35,
              size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
        cx += 2.1

# Connecting arrows between layers
for ay in [2.45, 3.55, 4.65, 5.75]:
    rect(s6, 6.55, ay, 0.1, 0.1, fill=LIGHTBLUE)
    txbox(s6, "▼", 6.44, ay - 0.04, 0.32, 0.22,
          size=10, color=LIGHTBLUE, align=PP_ALIGN.CENTER)

add_notes(s6, "Describe the layered architecture. UI layer sits on top, passes input to the "
              "agent layer. Agents use utilities. Data is persisted in CSV. Phase 2 components "
              "are shown at the bottom in a darker colour to indicate they are not yet active.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Five Agent Design
# ══════════════════════════════════════════════════════════════════════════════
s7 = add_slide()
rect(s7, 0, 0, 13.33, 7.5, fill=OFFWHITE)
header_bar(s7, "Five-Agent Architecture", "Modular design — each agent owns a single responsibility")
footer_bar(s7)

agents_data = [
    ("✅", BLUE,                     "SubmissionAgent",    "Accept description & log file\nValidate inputs\nGenerate BugReport object",
     "Text + log bytes",  "BugReport dataclass"),
    ("✅", RGBColor(0x0A,0x6A,0x3F), "PreprocessingAgent", "Lowercase, strip HTML\nRemove punctuation\nCollapse whitespace",
     "BugReport",         "Cleaned text dict"),
    ("🔒", RGBColor(0x4A,0x4A,0x7A), "EmbeddingAgent",     "Load SentenceTransformer\nEncode text → vector\n(Phase 2 — stub only)",
     "Cleaned text",      "Float vector list"),
    ("🔒", RGBColor(0x4A,0x4A,0x7A), "RetrievalAgent",     "Query ChromaDB\nReturn top-K matches\n(Phase 2 — stub only)",
     "Query vector",      "Similar bug list"),
    ("🔒", RGBColor(0x4A,0x4A,0x7A), "FixAdvisorAgent",    "Build RAG prompt\nCall Gemini API\n(Phase 2 — stub only)",
     "Bug + context",     "Fix advisory (MD)"),
]

# Header row
hcols = [0.35, 1.8, 4.4, 7.4, 10.0]
hwidths = [1.3, 2.45, 2.85, 2.45, 3.1]
hlabels = ["Status", "Agent", "Responsibilities", "Inputs", "Outputs"]
for hx, hw, hl in zip(hcols, hwidths, hlabels):
    rect(s7, hx, 1.55, hw, 0.48, fill=NAVY)
    txbox(s7, hl, hx + 0.06, 1.6, hw - 0.1, 0.38,
          size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

y = 2.15
for badge, col, name, resp, inp, out in agents_data:
    row_h = 0.88
    # Status badge
    rect(s7, 0.35, y, 1.3, row_h, fill=col)
    txbox(s7, badge, 0.35, y + 0.22, 1.3, 0.44,
          size=20, color=WHITE, align=PP_ALIGN.CENTER)
    # Agent name
    rect(s7, 1.8, y, 2.45, row_h, fill=WHITE, line=RGBColor(0xCC,0xDD,0xEE))
    txbox(s7, name, 1.88, y + 0.18, 2.3, 0.52,
          size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    # Responsibilities
    rect(s7, 4.4, y, 2.85, row_h, fill=OFFWHITE, line=RGBColor(0xCC,0xDD,0xEE))
    txbox(s7, resp, 4.5, y + 0.06, 2.65, 0.78, size=9, color=DARKGRAY)
    # Inputs
    rect(s7, 7.4, y, 2.45, row_h, fill=WHITE, line=RGBColor(0xCC,0xDD,0xEE))
    txbox(s7, inp, 7.5, y + 0.22, 2.25, 0.44, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    # Outputs
    rect(s7, 10.0, y, 3.1, row_h, fill=OFFWHITE, line=RGBColor(0xCC,0xDD,0xEE))
    txbox(s7, out, 10.08, y + 0.22, 2.92, 0.44, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    y += row_h + 0.04

add_notes(s7, "Explain the single-responsibility principle behind each agent. "
              "Phase 1 agents (green checkmarks) are fully working. "
              "Phase 2 agents (padlock) are implemented as stub classes that raise "
              "NotImplementedError — their interfaces are defined and ready.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Knowledge Base Design
# ══════════════════════════════════════════════════════════════════════════════
s8 = add_slide()
rect(s8, 0, 0, 13.33, 7.5, fill=OFFWHITE)
header_bar(s8, "Knowledge Base Design", "bug_dataset.csv — 14-field structured schema")
footer_bar(s8)

fields = [
    ("BugID",               "string",  "Auto-generated unique ID  (e.g. BUG-A3F2C1B0)"),
    ("Title",               "string",  "Short one-line summary (first 80 chars)"),
    ("Description",         "text",    "Full bug description entered by developer"),
    ("Module",              "string",  "Affected system module or service name"),
    ("Priority",            "enum",    "Critical / High / Medium / Low"),
    ("Severity",            "enum",    "Critical / High / Medium / Low"),
    ("ProgrammingLanguage", "string",  "e.g. Python, Java, JavaScript"),
    ("OperatingSystem",     "string",  "e.g. Linux, Windows, macOS"),
    ("StackTrace",          "text",    "Raw stack trace or uploaded log content"),
    ("ExpectedBehavior",    "text",    "What the developer expected to happen"),
    ("ActualBehavior",      "text",    "What actually happened"),
    ("FixStatus",           "enum",    "Open / In Progress / Fixed / Closed"),
    ("SubmittedAt",         "datetime","ISO-8601 UTC timestamp of submission"),
    ("LogFileName",         "string",  "Original filename of uploaded log (if any)"),
]

# Column headers
cols_x = [0.35, 3.5, 5.35, 7.35]
col_w  = [3.0,  1.7, 1.85, 5.6]
col_lbl= ["Field Name", "Type", "Key?", "Description"]
for cx, cw, cl in zip(cols_x, col_w, col_lbl):
    rect(s8, cx, 1.55, cw, 0.42, fill=NAVY)
    txbox(s8, cl, cx + 0.06, 1.6, cw - 0.1, 0.32,
          size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

key_fields = {"BugID", "Title", "Description", "FixStatus", "SubmittedAt"}
y = 2.05
row_h = 0.34
for i, (fname, ftype, fdesc) in enumerate(fields):
    bg = WHITE if i % 2 == 0 else OFFWHITE
    is_key = fname in key_fields
    # Row background
    rect(s8, 0.35, y, 12.6, row_h, fill=bg, line=RGBColor(0xCC,0xDD,0xEE))
    txbox(s8, fname,  0.45, y + 0.04, 2.9, 0.26, size=10, bold=is_key,
          color=BLUE if is_key else DARKGRAY)
    txbox(s8, ftype,  3.55, y + 0.04, 1.6, 0.26, size=9, color=GRAY, align=PP_ALIGN.CENTER)
    txbox(s8, "PK" if fname == "BugID" else ("✔" if is_key else ""),
          5.4, y + 0.04, 1.7, 0.26, size=9,
          color=GREEN if is_key else GRAY, align=PP_ALIGN.CENTER, bold=is_key)
    txbox(s8, fdesc,  7.4, y + 0.04, 5.5, 0.26, size=9, color=DARKGRAY)
    y += row_h

add_notes(s8, "Walk through the schema. Point out that BugID is the primary key. "
              "The StackTrace field stores uploaded log content. "
              "This schema is designed to map directly to ChromaDB metadata fields in Phase 2.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Bug Submission Module
# ══════════════════════════════════════════════════════════════════════════════
s9 = add_slide()
rect(s9, 0, 0, 13.33, 7.5, fill=OFFWHITE)
header_bar(s9, "Bug Submission Module", "Task 5 — Fully working Streamlit application")
footer_bar(s9)

# Left: UI mockup
rect(s9, 0.35, 1.6, 6.2, 5.2, fill=RGBColor(0x0D,0x1B,0x2E))
txbox(s9, "🐛  AI Smart Bug Analyzer", 0.55, 1.72, 5.8, 0.4,
      size=13, bold=True, color=WHITE)
rect(s9, 0.35, 2.18, 6.2, 0.03, fill=LIGHTBLUE)

# Tabs mock
for ti, (tlabel, tcol) in enumerate([
    ("✏️  Description", BLUE),
    ("📁  Upload Log",  NAVY),
    ("🔀  Both",        NAVY),
]):
    tx = 0.45 + ti * 2.02
    rect(s9, tx, 2.25, 1.92, 0.38, fill=tcol)
    txbox(s9, tlabel, tx + 0.05, 2.3, 1.82, 0.28, size=9, color=WHITE)

# Text area mock
rect(s9, 0.45, 2.72, 6.0, 1.3, fill=RGBColor(0x06,0x10,0x1E),
     line=RGBColor(0x33,0x55,0x88))
txbox(s9, "Paste your bug description here…\n\nNullPointerException in\nUserService.getUser() at line 42",
      0.55, 2.78, 5.8, 1.15, size=9, color=RGBColor(0x66,0x88,0xAA), italic=True)
txbox(s9, "142 / 10,000 characters", 4.0, 4.08, 2.3, 0.25, size=8,
      color=RGBColor(0x66,0xFF,0xAA), align=PP_ALIGN.RIGHT)

# Buttons mock
rect(s9, 0.45, 4.42, 2.3, 0.48, fill=BLUE)
txbox(s9, "🐛  Submit Bug", 0.5, 4.5, 2.2, 0.32, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s9, 2.85, 4.42, 1.3, 0.48, fill=NAVY)
txbox(s9, "🔄  Reset", 2.9, 4.5, 1.2, 0.32, size=10, color=WHITE, align=PP_ALIGN.CENTER)

# Success message mock
rect(s9, 0.45, 5.05, 6.0, 0.7, fill=RGBColor(0x0A,0x2A,0x12))
txbox(s9, "✅  Bug submitted!  ID: BUG-A3F2C1B0\n    Submitted: 2026-07-06 09:15:00 UTC  ·  Source: TEXT",
      0.55, 5.1, 5.8, 0.55, size=9, color=RGBColor(0x66,0xFF,0xAA))

# Right: feature list
rect(s9, 7.0, 1.6, 5.9, 5.2, fill=WHITE, line=RGBColor(0xBB,0xCC,0xDD))
txbox(s9, "Module Features", 7.15, 1.7, 5.5, 0.4, size=13, bold=True, color=NAVY)

features = [
    (GREEN, "✔  Paste bug description (text area)"),
    (GREEN, "✔  Upload .txt / .log files"),
    (GREEN, "✔  Character counter (live)"),
    (GREEN, "✔  File size & type validation"),
    (GREEN, "✔  Submit & Reset buttons"),
    (GREEN, "✔  Inline error messages"),
    (GREEN, "✔  Auto-generated Bug ID"),
    (GREEN, "✔  UTC timestamp"),
    (GREEN, "✔  CSV persistence"),
    (GREEN, "✔  Upload file saved to disk"),
    (GREEN, "✔  Cleaned text preview"),
    (BLUE,  "◉  Dataset Viewer tab"),
    (BLUE,  "◉  Documentation tab"),
]
fy = 2.2
for fcol, ftxt in features:
    txbox(s9, ftxt, 7.2, fy, 5.5, 0.32, size=10, color=DARKGRAY)
    fy += 0.33

add_notes(s9, "Demo the actual running application here if possible. "
              "The left panel shows a UI mockup of the submission form. "
              "Point out the character counter, the tab structure, the success message "
              "with the auto-generated Bug ID, and that the data is immediately visible "
              "in the Dataset Viewer tab.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Technologies Used
# ══════════════════════════════════════════════════════════════════════════════
s10 = add_slide()
rect(s10, 0, 0, 13.33, 7.5, fill=OFFWHITE)
header_bar(s10, "Technologies Used", "Phase 1 stack  ·  Phase 2 planned technologies")
footer_bar(s10)

# Phase 1 tech
rect(s10, 0.35, 1.6, 12.6, 0.4, fill=NAVY)
txbox(s10, "✅  Phase 1 — Active Technologies", 0.5, 1.66, 12, 0.28,
      size=12, bold=True, color=WHITE)

phase1_tech = [
    ("🐍", "Python 3.11",  BLUE,   "Core language. Type hints, dataclasses,\npathlib, csv, uuid, datetime."),
    ("🌊", "Streamlit",    BLUE,   "Web UI framework. Tabs, file uploader,\nmetrics, session state."),
    ("🐼", "Pandas",       RGBColor(0x0A,0x6A,0x3F), "CSV reading, filtering, and\ndataframe display in the viewer tab."),
    ("📄", "CSV / stdlib", RGBColor(0x8B,0x45,0x00), "Built-in csv module for persistent\nbug storage without a database."),
    ("💻", "VS Code / Kiro", RGBColor(0x4A,0x23,0x5A), "IDE with AI assistant for\ncode scaffolding and review."),
    ("🔀", "Git",          RGBColor(0x6A,0x10,0x10), "Version control. Branch strategy\nfor phase-by-phase development."),
]

x_positions2 = [0.35, 2.48, 4.61, 6.74, 8.87, 11.0]
w = 1.95
for i, (icon, name, col, desc) in enumerate(phase1_tech):
    x = x_positions2[i]
    rect(s10, x, 2.1, w, 3.0, fill=WHITE, line=RGBColor(0xCC,0xDD,0xEE))
    rect(s10, x, 2.1, w, 0.7, fill=col)
    txbox(s10, icon, x, 2.1, w, 0.7, size=24, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s10, name, x + 0.05, 2.88, w - 0.1, 0.38,
          size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    rect(s10, x + 0.1, 3.32, w - 0.2, 0.03, fill=RGBColor(0xDD,0xEE,0xFF))
    txbox(s10, desc, x + 0.1, 3.42, w - 0.2, 0.62, size=8.5, color=GRAY)

# Phase 2 section
rect(s10, 0.35, 5.3, 12.6, 0.38, fill=RGBColor(0x3A,0x3A,0x5A))
txbox(s10, "🔒  Phase 2 — Future Technologies  (not yet implemented)",
      0.5, 5.36, 12.2, 0.28, size=11, bold=True, color=RGBColor(0xAA,0xAA,0xCC))

future_tech = [
    ("ChromaDB",      "Local vector database for semantic bug storage"),
    ("SentenceTransformers", "Text-to-vector embedding model (all-MiniLM-L6-v2)"),
    ("Google Gemini", "LLM for AI-powered fix recommendation"),
    ("python-dotenv", "Secure API key management via .env"),
]
fx = 0.5
for fname, fdesc in future_tech:
    rect(s10, fx, 5.78, 2.8, 0.78, fill=RGBColor(0x1A,0x1A,0x2E),
         line=RGBColor(0x44,0x44,0x88))
    txbox(s10, f"🔒  {fname}", fx + 0.1, 5.85, 2.6, 0.3, size=10, bold=True,
          color=RGBColor(0x88,0x88,0xCC))
    txbox(s10, fdesc, fx + 0.1, 6.18, 2.6, 0.32, size=8.5,
          color=RGBColor(0x66,0x66,0x99))
    fx += 3.1

add_notes(s10, "For each Phase 1 technology, briefly explain why it was chosen. "
              "Python for rapid development, Streamlit to avoid writing HTML/JS, "
              "Pandas for zero-config CSV handling. Mention that the Phase 2 tools "
              "are already listed in requirements.txt as commented-out dependencies.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Current Progress
# ══════════════════════════════════════════════════════════════════════════════
s11 = add_slide()
rect(s11, 0, 0, 13.33, 7.5, fill=OFFWHITE)
header_bar(s11, "Current Progress", "Phase 1 complete  ·  Phase 2 roadmap defined")
footer_bar(s11)

# Progress bar visual
rect(s11, 0.35, 1.6, 12.6, 0.5, fill=RGBColor(0xDD,0xEE,0xFF))
rect(s11, 0.35, 1.6, 6.3, 0.5, fill=GREEN)
txbox(s11, "Phase 1 Complete — 50%", 0.5, 1.66, 5.8, 0.36,
      size=11, bold=True, color=WHITE)
txbox(s11, "Phase 2 Planned — 50%", 7.0, 1.66, 5.5, 0.36,
      size=11, bold=True, color=GRAY)

# Completed column
rect(s11, 0.35, 2.25, 5.85, 4.55, fill=WHITE, line=RGBColor(0xBB,0xCC,0xDD))
rect(s11, 0.35, 2.25, 5.85, 0.5, fill=GREEN)
txbox(s11, "✅  COMPLETED", 0.5, 2.32, 5.5, 0.36, size=13, bold=True, color=WHITE)

completed = [
    ("Task 1", "Core concepts documentation\n(RAG, embeddings, vector search explained)"),
    ("Task 2", "System architecture & data flow\n(4 Mermaid diagrams, component docs)"),
    ("Task 3", "Five agent classes\n(2 working + 3 stubs with full interfaces)"),
    ("Task 4", "Knowledge base schema\n(14-field CSV, JSON schema, Phase 2 DB design)"),
    ("Task 5", "Bug submission Streamlit app\n(validation, persistence, 3-tab UI)"),
]
cy = 2.9
for task, detail in completed:
    rect(s11, 0.45, cy, 1.0, 0.62, fill=GREEN)
    txbox(s11, task, 0.45, cy + 0.1, 1.0, 0.42, size=9, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s11, detail, 1.55, cy + 0.04, 4.45, 0.56, size=10, color=DARKGRAY)
    cy += 0.76

# Upcoming column
rect(s11, 7.1, 2.25, 5.85, 4.55, fill=WHITE, line=RGBColor(0xBB,0xCC,0xDD))
rect(s11, 7.1, 2.25, 5.85, 0.5, fill=RGBColor(0x3A,0x3A,0x5A))
txbox(s11, "🔒  UPCOMING — Phase 2", 7.25, 2.32, 5.5, 0.36,
      size=13, bold=True, color=WHITE)

upcoming = [
    ("Step 6", "EmbeddingAgent implementation\nSentenceTransformers — all-MiniLM-L6-v2"),
    ("Step 7", "ChromaDB vector store setup\nDataset ingestion pipeline"),
    ("Step 8", "RetrievalAgent — semantic search\nTop-K cosine similarity queries"),
    ("Step 9", "RAG pipeline integration\nGemini API connection"),
    ("Step 10", "Fix Advisory UI\nConfidence scores & similarity bars"),
]
uy = 2.9
for step, detail in upcoming:
    rect(s11, 7.2, uy, 1.0, 0.62, fill=RGBColor(0x3A,0x3A,0x5A))
    txbox(s11, step, 7.2, uy + 0.1, 1.0, 0.42, size=9, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s11, detail, 8.3, uy + 0.04, 4.45, 0.56, size=10, color=GRAY)
    uy += 0.76

add_notes(s11, "The progress bar shows 50% — Phase 1 is a complete, working system. "
              "Phase 2 will build on this foundation. Each upcoming step maps directly "
              "to one of the placeholder agent stubs already defined in the codebase.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
s12 = add_slide()
rect(s12, 0, 0, 13.33, 7.5, fill=OFFWHITE)
header_bar(s12, "Conclusion", "Phase 1 — A solid foundation for AI-powered bug analysis")
footer_bar(s12)

# Summary box
rect(s12, 0.35, 1.65, 8.2, 5.1, fill=WHITE, line=RGBColor(0xBB,0xCC,0xDD))
txbox(s12, "What we have built", 0.5, 1.75, 7.8, 0.4, size=14, bold=True, color=NAVY)

summary_pts = [
    "A professional multi-agent Python project with clean OOP architecture.",
    "A fully working bug submission UI with validation, persistence, and file uploads.",
    "A structured 14-field knowledge base with 20 pre-loaded sample bugs.",
    "Modular agent design that separates concerns and simplifies future extension.",
    "Detailed documentation covering RAG, embeddings, and vector search concepts.",
    "Placeholder agents that define the Phase 2 interfaces — ready to be filled in.",
]
sy = 2.28
for pt in summary_pts:
    txbox(s12, f"  ✔   {pt}", 0.5, sy, 7.9, 0.44, size=11, color=DARKGRAY)
    sy += 0.48

# Readiness callout
rect(s12, 0.35, 5.5, 8.2, 0.75, fill=BLUE)
txbox(s12, "🚀  The project is fully functional and ready for Phase 2 AI integration.",
      0.5, 5.58, 8.0, 0.55, size=12, bold=True, color=WHITE)

# Right key insight panel
rect(s12, 9.0, 1.65, 3.95, 5.1, fill=NAVY)
txbox(s12, "Key Insight", 9.15, 1.75, 3.6, 0.4, size=13, bold=True, color=LIGHTBLUE)
insight = (
    "By separating the submission "
    "layer from the AI layer, Phase 2 "
    "features can be added without "
    "changing any existing code.\n\n"
    "The agent interfaces are already "
    "defined — Phase 2 simply fills "
    "in the implementations."
)
txbox(s12, insight, 9.15, 2.28, 3.65, 3.0, size=11, color=WHITE)
txbox(s12, "Design Pattern:\nOpen/Closed Principle", 9.15, 5.55, 3.65, 0.65,
      size=10, bold=True, color=LIGHTBLUE)

add_notes(s12, "Summarise the key achievement: a fully working, well-structured system "
              "that is ready for AI enhancement. Emphasise the open/closed design — "
              "Phase 2 adds new behaviour without modifying existing code.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Future Scope
# ══════════════════════════════════════════════════════════════════════════════
s13 = add_slide()
rect(s13, 0, 0, 13.33, 7.5, fill=OFFWHITE)
header_bar(s13, "Future Scope", "Phase 2 and beyond — the AI enhancement roadmap")
footer_bar(s13)

future_items = [
    (BLUE,                        "🧠", "Text Embeddings",       "Integrate SentenceTransformers to convert bug descriptions into 384-dimension vectors for semantic comparison."),
    (RGBColor(0x0A,0x6A,0x3F),    "📦", "ChromaDB Integration",  "Set up a local ChromaDB persistent vector store. Ingest all 20 sample bugs and every new submission automatically."),
    (RGBColor(0x6A,0x0A,0x6A),    "🔍", "Semantic Search",       "Implement RetrievalAgent to find the top-K most similar past bugs using cosine similarity — no keyword matching needed."),
    (RGBColor(0x8B,0x45,0x00),    "🤖", "RAG + Gemini AI",       "Connect FixAdvisorAgent to Google Gemini. Use retrieved bugs as RAG context to generate precise, codebase-aware fixes."),
    (RGBColor(0x1A,0x4A,0x6A),    "🔗", "GitHub Integration",    "Auto-create GitHub issues from submitted bugs. Link fix commits back to bug IDs for full traceability."),
    (RGBColor(0x4A,0x1A,0x4A),    "📊", "Analytics Dashboard",   "Add charts showing bug severity trends, fix rates by module, and model confidence scores over time."),
]

positions = [(0.35, 1.6), (4.55, 1.6), (8.75, 1.6),
             (0.35, 4.2), (4.55, 4.2), (8.75, 4.2)]
for (x, y), (col, icon, title, desc) in zip(positions, future_items):
    rect(s13, x, y, 3.9, 2.35, fill=WHITE, line=RGBColor(0xCC,0xDD,0xEE))
    rect(s13, x, y, 3.9, 0.7, fill=col)
    txbox(s13, icon, x, y + 0.05, 0.7, 0.6, size=22, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s13, title, x + 0.72, y + 0.14, 3.0, 0.42, size=12, bold=True, color=WHITE)
    txbox(s13, desc, x + 0.1, y + 0.8, 3.7, 1.45, size=10, color=DARKGRAY)

add_notes(s13, "Each card in this slide maps to one of the placeholder agents already "
              "stubbed out in the codebase. Phases can be tackled independently since "
              "each agent has a defined interface. Mention the GitHub and analytics items "
              "as stretch goals beyond the core AI pipeline.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
s14 = add_slide()
rect(s14, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s14, 0, 5.0, 13.33, 0.06, fill=LIGHTBLUE)
rect(s14, 0, 5.06, 13.33, 2.44, fill=RGBColor(0x0A,0x12,0x2A))

# Bug icon
rect(s14, 5.9, 0.6, 1.5, 1.5, fill=BLUE)
txbox(s14, "🐛", 5.9, 0.6, 1.5, 1.5, size=50, color=WHITE, align=PP_ALIGN.CENTER)

txbox(s14, "Thank You", 0, 2.3, 13.33, 1.0,
      size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s14, "AI Smart Bug Analyzer & Fix Advisor", 0, 3.35, 13.33, 0.55,
      size=18, color=LIGHTBLUE, align=PP_ALIGN.CENTER)

# Questions row
rect(s14, 3.5, 4.1, 6.33, 0.7, fill=BLUE)
txbox(s14, "❓  Questions & Discussion", 3.5, 4.16, 6.33, 0.55,
      size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom info
info_items = [
    ("📁", "Phase 1 — Tasks 1–5 Complete"),
    ("🐍", "Python 3.11  ·  Streamlit  ·  Pandas"),
    ("🚀", "Phase 2 AI pipeline — coming next"),
]
ix = 1.3
for icon, text in info_items:
    txbox(s14, f"{icon}  {text}", ix, 5.3, 3.5, 0.5,
          size=11, color=LIGHTBLUE, align=PP_ALIGN.CENTER)
    ix += 3.7

txbox(s14, "Your College Name  ·  Department of Computer Science  ·  July 2026",
      0, 6.5, 13.33, 0.36, size=10, color=RGBColor(0x44,0x66,0x88),
      align=PP_ALIGN.CENTER)

add_notes(s14, "Closing slide. Invite questions. If time permits, offer a live demo "
              "of the Streamlit app — run 'streamlit run app.py' and show the submission "
              "form, character counter, and the Dataset Viewer tab.")

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out_path = "AI_Smart_Bug_Analyzer_Presentation.pptx"
prs.save(out_path)
print(f"✅  Presentation saved → {out_path}")
print(f"    Slides: {len(prs.slides)}")
